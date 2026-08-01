"""Local PPTX input: export slides to PNG and read speaker notes."""

from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
from typing import List, Optional

from images import cleanup_old_assets, resize_slide_png
from paths import TARGET_IMAGE_SIZE

try:
    from pptx import Presentation  # type: ignore[import-untyped]

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pypdfium2 as pdfium  # type: ignore[import-untyped]

    HAS_PDFIUM = True
except ImportError:
    HAS_PDFIUM = False


def is_pptx_path(source: str) -> bool:
    """Return True if source looks like a local .pptx file path."""
    value = (source or "").strip()
    if not value:
        return False
    if value.lower().endswith(".pptx") and os.path.isfile(value):
        return True
    return os.path.isfile(value) and value.lower().endswith(".pptx")


def get_pptx_slide_count(pptx_path: str) -> int:
    """Return the number of slides in a PPTX file."""
    if not HAS_PPTX:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")
    presentation = Presentation(pptx_path)  # type: ignore[possibly-unbound]
    return len(presentation.slides)


def get_pptx_title(pptx_path: str) -> str:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")
    presentation = Presentation(pptx_path)  # type: ignore[possibly-unbound]
    core_props = presentation.core_properties
    title = (core_props.title or "").strip()
    if title:
        return title
    return os.path.splitext(os.path.basename(pptx_path))[0] or "pptx_export"


def _find_soffice() -> str:
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/libreoffice",
        "/usr/bin/soffice",
    ]
    for candidate in candidates:
        if candidate and os.path.isfile(candidate):
            return candidate
    raise RuntimeError(
        "LibreOffice is required to render PPTX slides to PNG. "
        "Install LibreOffice and ensure `soffice` is on your PATH."
    )


def _libreoffice_convert_to_pdf(pptx_path: str, temp_dir: str) -> str:
    """Convert a PPTX to PDF via LibreOffice headless."""
    soffice = _find_soffice()
    abs_pptx = os.path.abspath(pptx_path)
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        temp_dir,
        abs_pptx,
    ]
    print(f"🖨️  Rendering PPTX to PDF with LibreOffice: {' '.join(cmd)}")
    result = subprocess.run(cmd, capture_output=True, text=True, check=False)
    if result.returncode != 0:
        stderr = (result.stderr or result.stdout or "").strip()
        raise RuntimeError(f"LibreOffice PDF conversion failed: {stderr or 'unknown error'}")

    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    expected_pdf = os.path.join(temp_dir, f"{base_name}.pdf")
    if os.path.isfile(expected_pdf):
        return expected_pdf

    generated = sorted(
        os.path.join(temp_dir, name)
        for name in os.listdir(temp_dir)
        if name.lower().endswith(".pdf")
    )
    if not generated:
        raise RuntimeError("LibreOffice did not produce a PDF from the PPTX")
    if len(generated) > 1:
        raise RuntimeError(
            f"LibreOffice produced multiple PDF files ({len(generated)}); expected one."
        )
    return generated[0]


def _pdf_pages_to_png(
    pdf_path: str,
    output_dir: str,
    *,
    page_numbers: Optional[List[int]] = None,
) -> List[str]:
    """Render selected PDF pages (1-based) to slide_XX.png files."""
    if not HAS_PDFIUM:
        raise RuntimeError(
            "pypdfium2 is required for PPTX slide rendering. Install with: pip install pypdfium2"
        )

    pdf = pdfium.PdfDocument(pdf_path)  # type: ignore[possibly-unbound]
    page_count = len(pdf)
    if page_count == 0:
        raise RuntimeError("PDF produced from PPTX has no pages")

    if page_numbers is None:
        page_numbers = list(range(1, page_count + 1))
    else:
        for page_num in page_numbers:
            if page_num < 1 or page_num > page_count:
                raise RuntimeError(
                    f"Requested slide {page_num} but PDF only has {page_count} page(s)"
                )

    target_width = TARGET_IMAGE_SIZE[0]
    final_paths: List[str] = []
    for out_index, page_num in enumerate(page_numbers, start=1):
        page = pdf[page_num - 1]
        width, _height = page.get_size()
        scale = target_width / width if width > 0 else 2.0
        bitmap = page.render(scale=scale)
        dest_path = os.path.join(output_dir, f"slide_{out_index:02d}.png")
        bitmap.to_pil().save(dest_path)
        final_paths.append(dest_path)

    return final_paths


def _libreoffice_png_paths(
    pptx_path: str,
    output_dir: str,
    *,
    only_slide: Optional[int] = None,
    expected_slide_count: Optional[int] = None,
) -> List[str]:
    """Convert a PPTX to one PNG per slide via LibreOffice PDF + pypdfium2."""
    os.makedirs(output_dir, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="synth_pptx_") as temp_dir:
        pdf_path = _libreoffice_convert_to_pdf(pptx_path, temp_dir)

        if only_slide is not None:
            page_numbers = [only_slide]
        else:
            page_numbers = None

        final_paths = _pdf_pages_to_png(pdf_path, output_dir, page_numbers=page_numbers)

        if only_slide is None and expected_slide_count is not None:
            if len(final_paths) != expected_slide_count:
                raise RuntimeError(
                    f"PPTX has {expected_slide_count} slide(s) but PDF export produced "
                    f"{len(final_paths)} page(s). The deck may be corrupt or unsupported."
                )

        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        print(f"Exported {len(final_paths)} slide image(s) from {base_name}")
        return final_paths


def _expected_slide_png_paths(
    output_dir: str,
    slide_count: int,
    *,
    only_slide: Optional[int] = None,
) -> List[str]:
    if only_slide is not None:
        return [os.path.join(output_dir, "slide_01.png")]
    return [
        os.path.join(output_dir, f"slide_{index:02d}.png")
        for index in range(1, slide_count + 1)
    ]


def export_pptx_slides_to_png(
    pptx_path: str,
    output_dir: str,
    only_slide: Optional[int] = None,
) -> int:
    """Export PPTX slides to PNG files. Returns number of images saved."""
    os.makedirs(output_dir, exist_ok=True)

    expected_slide_count = get_pptx_slide_count(pptx_path)
    if only_slide is not None:
        if only_slide < 1 or only_slide > expected_slide_count:
            raise ValueError(
                f"--only-slide must be between 1 and {expected_slide_count}"
            )

    expected_paths = _expected_slide_png_paths(
        output_dir,
        expected_slide_count,
        only_slide=only_slide,
    )
    if all(os.path.isfile(path) for path in expected_paths):
        print(
            f"\n📥 Using {len(expected_paths)} existing slide PNG(s) in {output_dir}"
        )
        return len(expected_paths)

    png_paths = _libreoffice_png_paths(
        pptx_path,
        output_dir,
        only_slide=only_slide,
        expected_slide_count=expected_slide_count if only_slide is None else None,
    )

    saved = 0
    for png_path in png_paths:
        resize_slide_png(png_path)
        saved += 1

    if only_slide is not None:
        # Renumber the single exported slide to slide_01.png for downstream consistency.
        expected = os.path.join(output_dir, "slide_01.png")
        source = png_paths[0] if png_paths else None
        if source and source != expected:
            if os.path.isfile(expected):
                os.remove(expected)
            os.rename(source, expected)
        saved = 1 if os.path.isfile(expected) else 0

    print(f"Exported {saved} slide image(s) to {output_dir}")
    return saved


def export_pptx_speaker_notes(
    pptx_path: str,
    output_dir: str,
    only_slide: Optional[int] = None,
) -> List[str]:
    """Read speaker notes from a PPTX file."""
    if not HAS_PPTX:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")

    os.makedirs(output_dir, exist_ok=True)
    cleanup_old_assets(output_dir, "_notes.txt")

    presentation = Presentation(pptx_path)  # type: ignore[possibly-unbound]
    all_notes: List[str] = []
    notes_lines: List[str] = []

    for slide_idx, slide in enumerate(presentation.slides, start=1):
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        all_notes.append(notes_text)

    export_notes = all_notes
    if only_slide is not None:
        if only_slide < 1 or only_slide > len(all_notes):
            raise ValueError(f"--only-slide must be between 1 and {len(all_notes)}")
        export_notes = [all_notes[only_slide - 1]]

    for file_idx, notes_text in enumerate(export_notes, start=1):
        slide_num = only_slide if only_slide is not None else file_idx
        note_path = os.path.join(output_dir, f"slide_{file_idx:02d}_notes.txt")
        with open(note_path, "w", encoding="utf-8") as note_file:
            note_file.write(notes_text)
        if notes_text:
            print(f"✅ Found notes for slide {slide_num} ({len(notes_text)} chars)")
        else:
            print(f"⚠️  No notes found for slide {slide_num}")
        notes_lines.append(f"=== Slide {slide_num} ===")
        notes_lines.append(notes_text or "(No notes)")
        notes_lines.append("")

    notes_file = os.path.join(output_dir, "notes_all.txt")
    with open(notes_file, "w", encoding="utf-8") as notes_file_handle:
        notes_file_handle.write("\n".join(notes_lines))
    print(f"\n✅ Exported all notes to: {notes_file}")
    return export_notes
